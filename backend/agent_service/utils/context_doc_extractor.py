"""
Context document extractor — reads bytes of PDF / DOCX / PPTX / TXT files and
returns plain text suitable for feeding into context_parser.parse_user_context().

Supported formats:
  .txt            — decoded as UTF-8 (fallback latin-1)
  .pdf            — text layer extracted page by page via pypdf
  .docx           — paragraphs + table cells via python-docx
  .pptx           — all text frames across all slides via python-pptx
  .doc / .ppt     — attempted via the DOCX/PPTX parsers; fails gracefully

Extraction is intentionally lenient: any page or shape that raises an error is
skipped so a single corrupted element never blocks the whole document.
"""
import io
import logging

_log = logging.getLogger(__name__)

# Max characters kept per document.
# 60 000 chars ≈ 15 000 tokens — covers a full 25-page BI specification PDF.
# The spec_reader and context_parser both use Claude Sonnet which has a 200K
# token input limit, so this is well within budget.
MAX_DOC_CHARS = 60_000


def _ext(filename: str) -> str:
    return filename.rsplit(".", 1)[-1].lower() if "." in filename else ""


def _extract_txt(data: bytes) -> str:
    for enc in ("utf-8", "utf-8-sig", "latin-1"):
        try:
            return data.decode(enc)
        except UnicodeDecodeError:
            continue
    return data.decode("latin-1", errors="replace")


def _extract_pdf(data: bytes) -> str:
    try:
        import pypdf  # noqa: PLC0415
    except ImportError:
        print("[context_doc_extractor] ⚠ pypdf not installed — PDF extraction unavailable", flush=True)
        return "[PDF extraction unavailable — install pypdf]"
    try:
        reader = pypdf.PdfReader(io.BytesIO(data))
        parts: list[str] = []
        for page in reader.pages:
            try:
                text = page.extract_text() or ""
                if text.strip():
                    parts.append(text)
            except Exception:
                pass
        result = "\n".join(parts)
        print(
            f"[context_doc_extractor] PDF extracted: pages={len(reader.pages)}"
            f"  text_chars={len(result)}"
            f"  first200={result[:200].replace(chr(10), ' ')!r}",
            flush=True,
        )
        return result
    except Exception as exc:
        print(f"[context_doc_extractor] ⚠ PDF extraction failed: {exc}", flush=True)
        _log.warning("PDF extraction failed: %s", exc)
        return ""


def _extract_docx(data: bytes) -> str:
    try:
        import docx  # python-docx  # noqa: PLC0415
    except ImportError:
        return "[DOCX extraction unavailable — install python-docx]"
    try:
        document = docx.Document(io.BytesIO(data))
        parts: list[str] = []
        for para in document.paragraphs:
            if para.text.strip():
                parts.append(para.text)
        # Also grab table cell text
        for table in document.tables:
            for row in table.rows:
                for cell in row.cells:
                    if cell.text.strip():
                        parts.append(cell.text)
        return "\n".join(parts)
    except Exception as exc:
        _log.warning("DOCX extraction failed: %s", exc)
        return ""


def _extract_pptx(data: bytes) -> str:
    try:
        from pptx import Presentation  # python-pptx  # noqa: PLC0415
    except ImportError:
        return "[PPTX extraction unavailable — install python-pptx]"
    try:
        prs = Presentation(io.BytesIO(data))
        parts: list[str] = []
        for slide in prs.slides:
            slide_texts: list[str] = []
            for shape in slide.shapes:
                try:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            t = para.text.strip()
                            if t:
                                slide_texts.append(t)
                except Exception:
                    pass
            if slide_texts:
                parts.append(" | ".join(slide_texts))
        return "\n".join(parts)
    except Exception as exc:
        _log.warning("PPTX extraction failed: %s", exc)
        return ""


def extract_text(filename: str, data: bytes) -> str:
    """
    Extract plain text from a document file.

    filename  — original filename, used only to determine the format.
    data      — raw file bytes.

    Returns extracted text truncated to MAX_DOC_CHARS (4 000).
    Returns empty string on failure (never raises).
    """
    ext = _ext(filename)
    try:
        if ext == "txt":
            text = _extract_txt(data)
        elif ext == "pdf":
            text = _extract_pdf(data)
        elif ext in ("docx", "doc"):
            text = _extract_docx(data)
        elif ext in ("pptx", "ppt"):
            text = _extract_pptx(data)
        else:
            _log.warning("context_doc_extractor: unsupported extension '%s' for '%s'", ext, filename)
            return ""
    except Exception as exc:
        _log.warning("context_doc_extractor: unexpected error for '%s': %s", filename, exc)
        return ""

    text = text.strip()
    if len(text) > MAX_DOC_CHARS:
        print(
            f"[context_doc_extractor] truncating '{filename}': {len(text)} → {MAX_DOC_CHARS} chars",
            flush=True,
        )
        text = text[:MAX_DOC_CHARS] + "…"
    if not text:
        print(f"[context_doc_extractor] ⚠ '{filename}' produced empty text", flush=True)
    else:
        print(
            f"[context_doc_extractor] '{filename}' → {len(text)} chars extracted",
            flush=True,
        )
    return text


def merge_context(typed_text: str, doc_extracts: list[tuple[str, str]]) -> str:
    """
    Combine user-typed context with text extracted from uploaded documents.

    typed_text    — text from the UI textarea (may be empty).
    doc_extracts  — list of (filename, extracted_text) tuples.

    Returns a single string ready for parse_user_context().
    """
    parts: list[str] = []
    if typed_text and typed_text.strip():
        parts.append(typed_text.strip())
    for filename, text in doc_extracts:
        if text and text.strip():
            parts.append(f"[From document: {filename}]\n{text.strip()}")
        else:
            print(f"[context_doc_extractor] merge_context: '{filename}' skipped (empty)", flush=True)
    merged = "\n\n".join(parts)
    print(
        f"[context_doc_extractor] merge_context: typed={len(typed_text or '')}c"
        f"  docs={len(doc_extracts)}  merged_total={len(merged)}c",
        flush=True,
    )
    return merged
